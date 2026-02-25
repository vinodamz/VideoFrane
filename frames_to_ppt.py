"""
frames_to_ppt.py
────────────────
Convert a folder of images into a PowerPoint presentation.
One image per slide, full-screen, widescreen 16:9.

Usage:
    python frames_to_ppt.py [images_folder] [output.pptx] [options]

Examples:
    python frames_to_ppt.py
    python frames_to_ppt.py frames output.pptx
    python frames_to_ppt.py frames output.pptx --title "My Video Frames"
    python frames_to_ppt.py frames output.pptx --captions --bg-color 1a1a2e
"""

import os
import sys
import argparse
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


SUPPORTED = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"}
SLIDE_W_IN = 13.33   # 16:9 widescreen width  (inches)
SLIDE_H_IN = 7.5     # 16:9 widescreen height (inches)


# ── CLI ───────────────────────────────────────────────────────────────────────
def parse_args():
    p = argparse.ArgumentParser(description="Images folder → PowerPoint presentation")
    p.add_argument("folder",  nargs="?", default="frames",
                   help="Folder containing images (default: frames)")
    p.add_argument("output",  nargs="?", default="output.pptx",
                   help="Output .pptx filename (default: output.pptx)")
    p.add_argument("--title",    default="",
                   help="Optional title slide text")
    p.add_argument("--captions", action="store_true",
                   help="Add filename as caption on each slide")
    p.add_argument("--bg-color", default="000000", metavar="RRGGBB",
                   help="Slide background colour in hex (default: 000000 = black)")
    return p.parse_args()


# ── Helpers ───────────────────────────────────────────────────────────────────
def collect_images(folder: str) -> list[str]:
    if not os.path.isdir(folder):
        print(f"[ERROR] Folder not found: {folder}")
        sys.exit(1)
    files = sorted(
        os.path.join(folder, f)
        for f in os.listdir(folder)
        if os.path.splitext(f)[1].lower() in SUPPORTED
    )
    if not files:
        print(f"[ERROR] No images found in: {folder}")
        sys.exit(1)
    return files


def hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


def set_slide_background(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def fit_image(img_w: int, img_h: int, slide_w: float, slide_h: float):
    """
    Return (left, top, width, height) in inches so the image fills
    the slide while preserving aspect ratio (letterbox / pillarbox).
    """
    img_ratio  = img_w / img_h
    slide_ratio = slide_w / slide_h

    if img_ratio > slide_ratio:          # wide image → fit width
        w = slide_w
        h = slide_w / img_ratio
    else:                                # tall image → fit height
        h = slide_h
        w = slide_h * img_ratio

    left = (slide_w - w) / 2
    top  = (slide_h - h) / 2
    return left, top, w, h


def add_title_slide(prs: Presentation, title_text: str, bg_color: RGBColor):
    slide_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, bg_color)

    txbox = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.33), Inches(2.5)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_image_slide(prs: Presentation, image_path: str,
                    bg_color: RGBColor, captions: bool):
    slide_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, bg_color)

    # Get natural image dimensions for aspect-ratio calc
    with PILImage.open(image_path) as im:
        img_w, img_h = im.size

    left, top, w, h = fit_image(img_w, img_h, SLIDE_W_IN, SLIDE_H_IN)
    slide.shapes.add_picture(
        image_path,
        Inches(left), Inches(top),
        Inches(w),    Inches(h)
    )

    # Optional caption bar at the bottom
    if captions:
        caption_h = 0.35
        txbox = slide.shapes.add_textbox(
            Inches(0), Inches(SLIDE_H_IN - caption_h),
            Inches(SLIDE_W_IN), Inches(caption_h)
        )
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = os.path.basename(image_path)
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    args = parse_args()
    files     = collect_images(args.folder)
    out_path  = args.output
    bg_color  = hex_to_rgb(args.bg_color)

    print(f"\n{'='*55}")
    print(f"  Images folder : {os.path.abspath(args.folder)}")
    print(f"  Total images  : {len(files)}")
    print(f"  Output file   : {os.path.abspath(out_path)}")
    print(f"  Background    : #{args.bg_color.upper()}")
    print(f"  Captions      : {'yes' if args.captions else 'no'}")
    print(f"{'='*55}\n")

    # Create presentation (widescreen 16:9)
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # Optional title slide
    if args.title:
        add_title_slide(prs, args.title, bg_color)
        print(f"  ✔ Title slide added: \"{args.title}\"")

    # One image per slide
    for i, path in enumerate(files, 1):
        print(f"  Adding slide {i:>4}/{len(files)} : {os.path.basename(path)}", end="\r")
        add_image_slide(prs, path, bg_color, args.captions)

    print(f"\n\n  Saving → {out_path} …")
    prs.save(out_path)
    print(f"\n[DONE] Presentation saved: {os.path.abspath(out_path)}")
    print(f"[DONE] {len(files) + (1 if args.title else 0)} slides total.")


if __name__ == "__main__":
    main()
